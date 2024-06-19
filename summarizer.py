
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.document_loaders import PyPDFLoader
from langchain.chains.summarize import load_summarize_chain
from langchain.chat_models import ChatOpenAI
from langchain.vectorstores import FAISS
from langchain.embeddings import OpenAIEmbeddings
from langchain.chains.summarize import load_summarize_chain
import numpy as np
from sklearn.cluster import KMeans
from langchain import OpenAI
from docx import Document as docx
from pptx import Presentation
import warnings
from langchain.prompts import PromptTemplate
from langchain.docstore.document import Document
openai_api_key='Insert Key Here'
warnings.filterwarnings("ignore", category=UserWarning, message="Importing OpenAI from langchain root module is no longer supported.")


#model and tokenizer loading

def load_pdf(file_path):
    loader =  PyPDFLoader(file_path)
    pages = loader.load_and_split()
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=3500, chunk_overlap=50)
    texts = text_splitter.split_documents(pages)
    final_texts = ""
    for text in texts:
        #print(text)
        final_texts = final_texts + text.page_content
    return final_texts
def load_word_document(file_path):
    doc = docx(file_path)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return "\n".join(text)

def load_powerpoint(file_path):
    presentation = Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

#file loader and preprocessing
def file_preprocessing(file):
    if file.endswith('.pdf'):
        return load_pdf(file)
    elif file.endswith(".docx"):
        return load_word_document(file)
    elif file.endswith(".pptx") or file.endswith(".ppt"):
        return load_powerpoint(file)
    else:
        raise Exception("Invalid format.")



def main():


    #filepath = input("Upload your PDF file: ")
    llm = OpenAI(temperature=0, openai_api_key=openai_api_key)
    filepath = '/Users/abdulrahmankhafagy/Documents/Personal_Projects/Python/Latest-Linguify/commentary.pdf' 
    text = file_preprocessing(filepath)
    
    
    tokens = llm.get_num_tokens(text)
    text_splitter = RecursiveCharacterTextSplitter(separators=["\n\n", "\n", "\t"], chunk_size=10000, chunk_overlap=3000)
    docs = text_splitter.create_documents([text])
    
    input_tokens = llm.get_num_tokens(docs[0].page_content)
    print(len(docs))
    print(input_tokens)
    embeddings = OpenAIEmbeddings(openai_api_key=openai_api_key)
    vectors = embeddings.embed_documents([x.page_content for x in docs])
    # while True:
    #     try:
    num_clusters = len(docs)
    kmeans = KMeans(n_clusters=num_clusters, random_state=42, n_init=10).fit(vectors)
        #     break
        # except:
        #     num_clusters -= 2

    closest_indices = []

    # Loop through the number of clusters you have
    for i in range(num_clusters):
        
        # Get the list of distances from that particular cluster center
        distances = np.linalg.norm(vectors - kmeans.cluster_centers_[i], axis=1)
        
        # Find the list position of the closest one (using argmin to find the smallest distance)
        closest_index = np.argmin(distances)
        
        # Append that position to your closest indices list
        closest_indices.append(closest_index)
    selected_indices = sorted(closest_indices)
    llm3 = ChatOpenAI(temperature=0,
                 openai_api_key=openai_api_key,
                 max_tokens=1000,
                 model='gpt-3.5-turbo-16k'
                )
    map_prompt = """
    You will be given a passage of a book. This section will be enclosed in triple backticks (```)
    Your goal is to give a summary of this section so that a reader will have a full understanding of what happened.

    ```{text}```
    FULL SUMMARY:
    """
    map_prompt_template = PromptTemplate(template=map_prompt, input_variables=["text"])

    map_chain = load_summarize_chain(llm=llm3,
                             chain_type="stuff",
                             prompt=map_prompt_template)
    selected_docs = [docs[doc] for doc in selected_indices]

    summary_list = []

# Loop through a range of the lenght of your selected docs
    for i, doc in enumerate(selected_docs):
        
        # Go get a summary of the chunk
        chunk_summary = map_chain.run([doc])
        
        # Append that summary to your list
        summary_list.append(chunk_summary)
    summaries = "\n".join(summary_list)

    # Convert it back to a document
    summaries = Document(page_content=summaries)
    llm4 = ChatOpenAI(temperature=0,
                 openai_api_key=openai_api_key,
                 max_tokens=3000,
                 model='gpt-3.5-turbo-16k',
                 request_timeout=120
                )
    combine_prompt = """
    You are an assistant that only speaks in Markdown. Write text that is only formatted as markdown. The summaries will be enclosed in triple backticks (```)
    Write a Title for the notes that is under 15 words.\nThen write: "- Summary -"\nWrite "Summary" as a Heading 1.\nWrite a summary of the provided notes.\nThen write: "-Additional Info-_".\nThen you must return a list of the main points in the provided transcript. Then return a list of Key Concepts. Then return a list of Definitions and examples. Then return a list of Equations and Formulas if there are any. Then return a list of historical contexts and important dates if there are any.\nFor each list, return a Heading 2 before writing the list items. Limit each list item to 100 words, and return at least 5 points per list.

    ```{text}```
    SUMMARY:
    """
    combine_prompt_template = PromptTemplate(template=combine_prompt, input_variables=["text"])
    reduce_chain = load_summarize_chain(llm=llm4,
                                chain_type="stuff",
                                prompt=combine_prompt_template,
                                #verbose=True # Set this to true if you want to see the inner workings
                                    )
    
    output = reduce_chain.run([summaries])
    tokenss = llm.get_num_tokens(output)
    print (output)
    



if __name__ == "__main__":
    main()
